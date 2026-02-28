#!/usr/bin/env python3
"""
CSV/XLSX アンケート結果 → PowerPoint 変換ツール

- 個人情報（氏名・メールアドレス等）はPPTXに出力しません
- グラフはPowerPoint内でネイティブ編集可能な形式で生成します
- CSV / XLSX 両対応

使い方:
    python survey_to_pptx.py input.xlsx
    python survey_to_pptx.py input.csv -o output.pptx --title "顧客満足度調査"
"""

import argparse
import re
import sys
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 本文フォント（PPTX全体で統一）
FONT_NAME = "游ゴシック"
# 表のタイトル行（ヘッダー行）の背景色
TABLE_HEADER_COLOR = RGBColor(0x19, 0x20, 0x48)  # #192048

# ── テーマ・定数 ──────────────────────────────────────────────────────────────

THEME = {
    "primary":    RGBColor(0x1F, 0x4E, 0x79),
    "accent":     RGBColor(0x2E, 0x75, 0xB6),
    "accent2":    RGBColor(0x70, 0xAD, 0x47),
    "light":      RGBColor(0xDE, 0xEB, 0xF7),
    "text_dark":  RGBColor(0x26, 0x26, 0x26),
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),
    "separator":  RGBColor(0xE0, 0xE0, 0xE0),
}

CHART_COLORS = [
    RGBColor(0x2E, 0x75, 0xB6),
    RGBColor(0x70, 0xAD, 0x47),
    RGBColor(0xFF, 0xC0, 0x00),
    RGBColor(0xFF, 0x6B, 0x6B),
    RGBColor(0x9B, 0x59, 0xB6),
    RGBColor(0x1A, 0xBC, 0x9C),
    RGBColor(0xE6, 0x7E, 0x22),
    RGBColor(0x34, 0x98, 0xDB),
    RGBColor(0x2E, 0xCC, 0x71),
    RGBColor(0xE7, 0x4C, 0x3C),
]

# 4段階評価: 固定表示順と色
SCALE_4_ORDER  = ["大変参考になった", "参考になった", "どちらでもない", "参考にならなかった"]
SCALE_4_COLORS = [
    RGBColor(0x1F, 0x4E, 0x79),  # 濃青
    RGBColor(0x2E, 0x75, 0xB6),  # 青
    RGBColor(0xFF, 0xC0, 0x00),  # 黄
    RGBColor(0xFF, 0x6B, 0x6B),  # 赤
]

# テンプレート（既存6枚はサンプルとして残す）
DEFAULT_TEMPLATE_PATH = Path(__file__).resolve().parent / "template" / "design.pptx"

# テンプレートのレイアウト番号（スライドマスタの名前と対応）
LAYOUT_TITLE = 0          # タイトル
LAYOUT_SUMMARY = 1        # サマリ
LAYOUT_CHART = 2          # アンケート結果（グラフ）
LAYOUT_TABLE = 3          # アンケート結果（表）
LAYOUT_SECTION = 4        # セクションタイトル
LAYOUT_BODY = 5           # 本文

# テンプレート未使用時用のデフォルトサイズ
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

TOP_N_HIGH_CARDINALITY = 10  # 高カーディナリティ列の最大表示件数
# 表の「No.」列の幅（インチ）。回答列が残りの幅になる
TABLE_NO_COL_WIDTH_INCHES = 0.55
# 元データのこの列（ExcelでP列＝1始まりで16列目、0始まりで15）は円グラフ1枚で表示
PIE_ONLY_COLUMN_INDEX = 15

# ── 個人情報・メタデータ判定 ──────────────────────────────────────────────────

# 個人情報として除外する列名（小文字・部分一致）
_PERSONAL_EXACT = {
    "email", "e-mail", "mail", "メール", "メールアドレス",
    "name", "名前", "氏名", "姓", "名", "苗字", "名字",
    "contact last name", "contact first name",
    "last name", "first name", "lastname", "firstname",
    "record id", "contact id",
    "phone", "tel", "telephone", "電話", "電話番号", "携帯", "携帯番号", "mobile", "fax", "ファックス",
}
_PERSONAL_PATTERNS = [
    re.compile(r, re.IGNORECASE) for r in [
        r"^e-?mail", r"mail$",
        r"^(last|first)\s*name$",
        r"^(姓|名|氏名|苗字|名字)$",
        r"record.?id", r"contact.?id",
        r"phone", r"tel(ephone)?", r"電話", r"携帯", r"mobile", r"fax",
    ]
]
_EMAIL_RE = re.compile(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}")
# 電話番号らしい値（数字・ハイフン・括弧のみで 10 桁前後）
_PHONE_RE = re.compile(r"^[\d\-\(\)\s+]{10,}$")

# 集計せずタイトル生成にのみ使うメタデータ列
_METADATA_EXACT = {"date", "survey name", "survey_name"}

# Appendix として最後に表で出す列（会社名など）
_APPENDIX_TABLE_COL_NAMES = {"会社名", "company name"}


def _get_placeholder_by_idx(slide_or_layout, idx: int):
    """スライドまたはレイアウトから idx のプレースホルダーを取得"""
    for ph in slide_or_layout.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _remove_shape(shape) -> None:
    """スライドからシェイプを削除（空プレースホルダ用）"""
    try:
        sp = shape._element
        sp.getparent().remove(sp)
    except Exception:
        pass


def _apply_body_font(text_frame, bold: bool = False, color: RGBColor | None = None) -> None:
    """テキストフレーム内の全 Run にフォント（游ゴシック）を適用"""
    for p in text_frame.paragraphs:
        for run in p.runs:
            run.font.name = FONT_NAME
            if bold:
                run.font.bold = True
            if color is not None:
                run.font.color.rgb = color


def is_personal_col(col: str) -> bool:
    c = col.lower().strip()
    if c in _PERSONAL_EXACT:
        return True
    return any(p.search(c) for p in _PERSONAL_PATTERNS)


def is_metadata_col(col: str) -> bool:
    return col.lower().strip() in _METADATA_EXACT


def is_appendix_company_col(col: str) -> bool:
    """会社名など、Appendix で表として扱う列か"""
    return col.strip().lower() in _APPENDIX_TABLE_COL_NAMES


def has_email_values(series: pd.Series) -> bool:
    return series.dropna().astype(str).str.contains(_EMAIL_RE).any()


def has_phone_values(series: pd.Series) -> bool:
    """値が電話番号形式（数字・ハイフン等のみで10桁以上）のセルが一定割合あれば True"""
    s = series.dropna().astype(str).str.strip()
    s = s[(s != "") & (s != "nan")]
    if len(s) < 3:
        return False
    matches = s.str.match(_PHONE_RE)
    return matches.sum() / len(s) >= 0.5


# ── データ読み込み ────────────────────────────────────────────────────────────

def load_data(path: str, encoding: str = "utf-8") -> pd.DataFrame:
    p = Path(path)
    if p.suffix.lower() in (".xlsx", ".xls"):
        return pd.read_excel(path)
    for enc in [encoding, "utf-8-sig", "shift_jis", "cp932"]:
        try:
            return pd.read_csv(path, encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError(f"CSVのエンコーディングを判定できませんでした: {path}")


# ── 列タイプ判定 ──────────────────────────────────────────────────────────────

def detect_col_type(series: pd.Series) -> str:
    """personal / metadata / categorical / high_cardinality / text / empty"""
    s = series.dropna().astype(str).str.strip()
    s = s[(s != "") & (s != "nan") & (s != "N/A")]
    if len(s) == 0:
        return "empty"
    if s.str.len().mean() > 25:
        return "text"
    unique_ratio = s.nunique() / len(s)
    if s.nunique() <= 25 and unique_ratio < 0.6:
        return "categorical"
    if s.nunique() > 25:
        return "high_cardinality"
    return "text"


def classify_columns(df: pd.DataFrame) -> dict[str, str]:
    result = {}
    for col in df.columns:
        if is_personal_col(col) or has_email_values(df[col]) or has_phone_values(df[col]):
            result[col] = "personal"
        elif is_metadata_col(col):
            result[col] = "metadata"
        elif is_appendix_company_col(col):
            result[col] = "appendix_company"
        else:
            result[col] = detect_col_type(df[col])
    return result


# ── スライド共通ユーティリティ ────────────────────────────────────────────────

def _fill_slide(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height,
              fill_color: RGBColor, line_color: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def _add_textbox(slide, left, top, width, height,
                 text: str, font_size: int,
                 bold: bool = False,
                 color: RGBColor | None = None,
                 align: PP_ALIGN = PP_ALIGN.LEFT,
                 word_wrap: bool = True) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.font.color.rgb = color or THEME["text_dark"]
    p.alignment = align


def _add_slide_header(slide, question_label: str, question_num: int) -> None:
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), THEME["primary"])
    # Q番号バッジ
    _add_rect(slide, Inches(0.25), Inches(0.15), Inches(0.65), Inches(0.75),
              THEME["accent"])
    _add_textbox(slide, Inches(0.25), Inches(0.15), Inches(0.65), Inches(0.75),
                 f"Q{question_num}", 16, bold=True,
                 color=THEME["text_light"], align=PP_ALIGN.CENTER)
    # 質問テキスト
    _add_textbox(slide, Inches(1.05), Inches(0.15), Inches(11.9), Inches(0.75),
                 question_label, 17, bold=True,
                 color=THEME["text_light"], align=PP_ALIGN.LEFT)


def _add_response_count(slide, n_answered: int, total: int) -> None:
    _add_textbox(slide, Inches(9.5), Inches(6.9), Inches(3.5), Inches(0.45),
                 f"回答数: {n_answered} / {total} 名",
                 11, color=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.RIGHT)


# ── ネイティブ PPTX チャート ──────────────────────────────────────────────────

def _color_chart_points(chart, colors: list[RGBColor]) -> None:
    """各ポイント（バー・パイスライス）に個別の色を設定"""
    try:
        for i, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[i % len(colors)]
    except Exception:
        pass


def add_bar_chart(slide, counts: pd.Series,
                  left, top, width, height,
                  colors: list[RGBColor],
                  title: str = "") -> None:
    """水平棒グラフ（PowerPoint ネイティブ・編集可能）"""
    total = int(counts.sum())
    # カテゴリラベルにパーセンテージを含める
    cat_labels = [
        f"{str(cat)[:28]}  ({v / total * 100:.1f}%)"
        for cat, v in zip(counts.index, counts.values)
    ]
    chart_data = ChartData()
    chart_data.categories = cat_labels
    chart_data.add_series("回答数", tuple(int(v) for v in counts.values))

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = cf.chart
    chart.has_legend = False

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
    else:
        chart.has_title = False

    # データラベル（件数表示）
    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.show_value = True
    dls.font.name = FONT_NAME
    dls.font.size = Pt(10)
    dls.font.bold = True

    # 軸フォント
    chart.value_axis.tick_labels.font.name = FONT_NAME
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = FONT_NAME
    chart.category_axis.tick_labels.font.size = Pt(10)

    _color_chart_points(chart, colors)


def add_pie_chart(slide, counts: pd.Series,
                  left, top, width, height,
                  colors: list[RGBColor],
                  title: str = "") -> None:
    """円グラフ（PowerPoint ネイティブ・編集可能）"""
    cat_labels = [str(cat)[:20] for cat in counts.index]
    chart_data = ChartData()
    chart_data.categories = cat_labels
    chart_data.add_series("割合", tuple(int(v) for v in counts.values))

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    )
    chart = cf.chart

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
    else:
        chart.has_title = False

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.name = FONT_NAME
    chart.legend.font.size = Pt(10)

    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.show_percentage = True
    dls.show_value = False
    dls.show_category_name = False
    dls.font.name = FONT_NAME
    dls.font.size = Pt(11)
    dls.font.bold = True

    _color_chart_points(chart, colors)


# ── スライドビルダー ──────────────────────────────────────────────────────────

def add_title_slide(prs: Presentation, title: str, subtitle: str,
                    n_respondents: int) -> None:
    """テンプレートの「タイトル」レイアウトを使用"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    ph10 = _get_placeholder_by_idx(slide, 10)
    ph11 = _get_placeholder_by_idx(slide, 11)
    ph12 = _get_placeholder_by_idx(slide, 12)
    if ph10 and ph10.text_frame:
        ph10.text_frame.text = title
        _apply_body_font(ph10.text_frame)
    if ph11 and ph11.text_frame:
        lines = [subtitle] if subtitle else []
        lines.append(f"回答者数: {n_respondents} 名")
        ph11.text_frame.text = "\n".join(lines)
        _apply_body_font(ph11.text_frame)
    if ph12 and ph12.text_frame:
        ph12.text_frame.text = datetime.now().strftime("%Y年%m月%d日")
        _apply_body_font(ph12.text_frame)


def add_categorical_slide(prs: Presentation,
                          series: pd.Series,
                          question_label: str,
                          question_num: int,
                          total_respondents: int,
                          ordered_values: list | None = None,
                          colors: list[RGBColor] | None = None,
                          top_n: int | None = None) -> None:
    """テンプレートの「アンケート結果（グラフ）」レイアウト。常に左：棒グラフ・右：円グラフ（同一データ）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CHART])
    ph0 = _get_placeholder_by_idx(slide, 0)
    ph10 = _get_placeholder_by_idx(slide, 10)
    ph11 = _get_placeholder_by_idx(slide, 11)
    ph12 = _get_placeholder_by_idx(slide, 12)

    if ph0 and ph0.text_frame:
        ph0.text_frame.text = f"Q{question_num}  {question_label}"
        _apply_body_font(ph0.text_frame)

    s = series.dropna().astype(str).str.strip()
    s = s[(s != "") & (s != "nan") & (s != "N/A")]
    counts = s.value_counts()

    if ordered_values:
        ordered = [v for v in ordered_values if v in counts.index]
        rest    = [v for v in counts.index if v not in ordered_values]
        counts  = counts.reindex(ordered + rest).dropna()

    if top_n and len(counts) > top_n:
        others = counts.iloc[top_n:].sum()
        counts = counts.head(top_n)
        if others > 0:
            counts = pd.concat([counts, pd.Series({"その他": int(others)})])

    n_answered = int(s.count())
    chart_colors = colors or CHART_COLORS

    # 常に左：棒グラフ、右：円グラフ（同一データ）
    if ph10 and ph11:
        add_bar_chart(slide, counts, ph10.left, ph10.top, ph10.width, ph10.height,
                      colors=chart_colors, title="回答数")
        add_pie_chart(slide, counts, ph11.left, ph11.top, ph11.width, ph11.height,
                      colors=chart_colors, title="割合")
        _remove_shape(ph10)
        _remove_shape(ph11)
    else:
        add_bar_chart(slide, counts, Inches(0.3), Inches(1.2), Inches(5), Inches(5),
                      colors=chart_colors, title="回答数")
        add_pie_chart(slide, counts, Inches(5.5), Inches(1.2), Inches(5), Inches(5),
                      colors=chart_colors, title="割合")

    if ph12 and ph12.text_frame:
        ph12.text_frame.text = f"回答数: {n_answered} / {total_respondents} 名"
        _apply_body_font(ph12.text_frame)


def add_pie_only_slide(prs: Presentation,
                       series: pd.Series,
                       question_label: str,
                       question_num: int,
                       total_respondents: int,
                       colors: list[RGBColor] | None = None) -> None:
    """指定列（例: P列）を円グラフ1枚で表示。左・右プレースホルダーを合わせた領域に1つの円グラフを配置。"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CHART])
    ph0 = _get_placeholder_by_idx(slide, 0)
    ph10 = _get_placeholder_by_idx(slide, 10)
    ph11 = _get_placeholder_by_idx(slide, 11)
    ph12 = _get_placeholder_by_idx(slide, 12)

    if ph0 and ph0.text_frame:
        ph0.text_frame.text = f"Q{question_num}  {question_label}"
        _apply_body_font(ph0.text_frame)

    s = series.dropna().astype(str).str.strip()
    s = s[(s != "") & (s != "nan") & (s != "N/A")]
    counts = s.value_counts()
    n_answered = int(s.count())
    chart_colors = colors or CHART_COLORS

    if ph10 and ph11:
        # 左・右のプレースホルダーを合わせた領域に1つの円グラフを配置
        left = ph10.left
        top = min(ph10.top, ph11.top)
        width = (ph11.left + ph11.width) - ph10.left
        height = max(ph10.height, ph11.height)
        add_pie_chart(slide, counts, left, top, width, height,
                      colors=chart_colors, title="")
        _remove_shape(ph10)
        _remove_shape(ph11)
    else:
        add_pie_chart(slide, counts, Inches(0.3), Inches(1.2), Inches(10), Inches(5.2),
                      colors=chart_colors, title="")

    if ph12 and ph12.text_frame:
        ph12.text_frame.text = f"回答数: {n_answered} / {total_respondents} 名"
        _apply_body_font(ph12.text_frame)


def _style_table_cell(cell, font_size: int = 10, bold: bool = False,
                      color: RGBColor | None = None) -> None:
    """表のセルにフォント（游ゴシック）・サイズ・太字・色・上下中央揃えを設定"""
    c = color or THEME["text_dark"]
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = c


def add_text_slides(prs: Presentation,
                    series: pd.Series,
                    question_label: str,
                    question_num: int,
                    total_respondents: int,
                    items_per_slide: int = 10) -> None:
    """テンプレートの「アンケート結果（表）」レイアウトで自由記述を表形式に"""
    texts = (
        series.dropna().astype(str).str.strip()
        .pipe(lambda s: s[(s != "") & (s != "nan") & (s != "N/A")])
        .reset_index(drop=True)
    )
    if len(texts) == 0:
        return

    total_pages = (len(texts) - 1) // items_per_slide + 1
    n_cols = 2

    for page in range(total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TABLE])
        header_label = question_label + (f"  ({page + 1}/{total_pages})" if total_pages > 1 else "")

        ph0 = _get_placeholder_by_idx(slide, 0)
        if ph0 and ph0.text_frame:
            ph0.text_frame.text = f"Q{question_num}  {header_label}"
            _apply_body_font(ph0.text_frame)

        start = page * items_per_slide
        page_texts = texts[start: start + items_per_slide]
        n_rows = 1 + len(page_texts)

        table_ph = _get_placeholder_by_idx(slide, 10)
        table_shape = None
        if table_ph and hasattr(table_ph, "insert_table"):
            table_shape = table_ph.insert_table(n_rows, n_cols)
            table = table_shape.table
        else:
            table_left, table_top = Inches(0.4), Inches(1.15)
            table_width = prs.slide_width - int(Inches(0.8))
            table_height = Inches(0.5 + len(page_texts) * 0.45)
            table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
            table = table_shape.table

        no_col_width_emu = int(Inches(TABLE_NO_COL_WIDTH_INCHES))
        total_w = table_shape.width if table_shape else prs.slide_width
        table.columns[0].width = no_col_width_emu
        table.columns[1].width = total_w - no_col_width_emu

        table.cell(0, 0).text = "No."
        table.cell(0, 1).text = "回答"
        for c in range(n_cols):
            _style_table_cell(table.cell(0, c), font_size=11, bold=True, color=THEME["text_light"])
            table.cell(0, c).fill.solid()
            table.cell(0, c).fill.fore_color.rgb = TABLE_HEADER_COLOR

        for i, text in enumerate(page_texts):
            row = 1 + i
            table.cell(row, 0).text = str(start + i + 1)
            display = text if len(text) <= 200 else text[:197] + "..."
            table.cell(row, 1).text = display
            for c in range(n_cols):
                _style_table_cell(table.cell(row, c), font_size=10)

        ph12 = _get_placeholder_by_idx(slide, 12)
        if ph12 and ph12.text_frame:
            ph12.text_frame.text = f"回答数: {len(texts)} / {total_respondents} 名"
            _apply_body_font(ph12.text_frame)


def add_section_title_slide(prs: Presentation, title: str) -> None:
    """テンプレートの「セクションタイトル」レイアウトを使用（白文字・太字）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    ph0 = _get_placeholder_by_idx(slide, 0)
    if ph0 and ph0.text_frame:
        ph0.text_frame.text = title
        _apply_body_font(ph0.text_frame, bold=True, color=THEME["text_light"])


def add_appendix_company_slides(prs: Presentation,
                                series: pd.Series,
                                column_label: str,
                                total_respondents: int,
                                items_per_slide: int = 12) -> None:
    """テンプレートの「アンケート結果（表）」レイアウトで Appendix 会社名を表形式に"""
    values = (
        series.dropna().astype(str).str.strip()
        .pipe(lambda s: s[(s != "") & (s != "nan") & (s != "N/A")])
        .reset_index(drop=True)
    )
    if len(values) == 0:
        return

    total_pages = (len(values) - 1) // items_per_slide + 1
    n_cols = 2

    for page in range(total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TABLE])
        header_label = f"Appendix: {column_label}"
        if total_pages > 1:
            header_label += f"  ({page + 1}/{total_pages})"

        ph0 = _get_placeholder_by_idx(slide, 0)
        if ph0 and ph0.text_frame:
            ph0.text_frame.text = header_label
            _apply_body_font(ph0.text_frame)

        start = page * items_per_slide
        page_values = values[start: start + items_per_slide]
        n_rows = 1 + len(page_values)

        table_ph = _get_placeholder_by_idx(slide, 10)
        if table_ph and hasattr(table_ph, "insert_table"):
            table_shape = table_ph.insert_table(n_rows, n_cols)
            table = table_shape.table
        else:
            table_left, table_top = Inches(0.4), Inches(1.15)
            table_width = prs.slide_width - int(Inches(0.8))
            table_height = Inches(0.5 + len(page_values) * 0.4)
            table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
            table = table_shape.table

        no_col_width_emu = int(Inches(TABLE_NO_COL_WIDTH_INCHES))
        total_w = table_shape.width if table_shape else prs.slide_width
        table.columns[0].width = no_col_width_emu
        table.columns[1].width = total_w - no_col_width_emu

        table.cell(0, 0).text = "No."
        table.cell(0, 1).text = column_label
        for c in range(n_cols):
            _style_table_cell(table.cell(0, c), font_size=11, bold=True, color=THEME["text_light"])
            table.cell(0, c).fill.solid()
            table.cell(0, c).fill.fore_color.rgb = TABLE_HEADER_COLOR

        for i, val in enumerate(page_values):
            row = 1 + i
            table.cell(row, 0).text = str(start + i + 1)
            table.cell(row, 1).text = val[:200] if len(val) <= 200 else val[:197] + "..."
            for c in range(n_cols):
                _style_table_cell(table.cell(row, c), font_size=10)

        ph12 = _get_placeholder_by_idx(slide, 12)
        if ph12 and ph12.text_frame:
            ph12.text_frame.text = f"回答数: {len(values)} / {total_respondents} 名"
            _apply_body_font(ph12.text_frame)


def add_summary_slide(prs: Presentation, df: pd.DataFrame,
                      col_types: dict, total_respondents: int) -> None:
    """テンプレートの「本文」レイアウトを使用し、サマリーカードを配置"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BODY])
    ph0 = _get_placeholder_by_idx(slide, 0)
    if ph0 and ph0.text_frame:
        ph0.text_frame.text = "サマリー"
        _apply_body_font(ph0.text_frame)

    _add_stat_card(slide, "総回答者数", f"{total_respondents} 名",
                   Inches(0.4), Inches(1.2), THEME["primary"])

    graph_cols = [
        (col, t) for col, t in col_types.items()
        if t in ("categorical", "high_cardinality")
    ][:6]

    # テンプレート幅に合わせて6枚のカードを2段×3列で配置
    card_w = Inches(2.35)
    card_gap = Inches(0.15)
    for i, (col_name, _) in enumerate(graph_cols):
        x = Inches(3.35) + (i % 3) * (card_w + card_gap)
        y = Inches(1.2) if i < 3 else Inches(2.95)
        s = df[col_name].dropna().astype(str).str.strip()
        s = s[(s != "") & (s != "nan") & (s != "N/A")]
        vc = s.value_counts()
        if len(vc):
            top = str(vc.index[0])[:20]
            pct = vc.iloc[0] / vc.sum() * 100
            value_text = f"最多: {top}\n({pct:.0f}%)"
        else:
            value_text = "データなし"
        _add_stat_card(slide, str(col_name)[:25], value_text,
                       x, y, THEME["accent"], width=card_w, height=Inches(1.55))


def _add_stat_card(slide, title: str, value: str,
                   left, top, color: RGBColor,
                   width=None, height=None) -> None:
    cw = width or Inches(3.7)
    ch = height or Inches(1.55)
    _add_rect(slide, left, top, cw, ch, color)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                 cw - Inches(0.3), Inches(0.55),
                 title, 12, color=THEME["text_light"])
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.68),
                 cw - Inches(0.3), Inches(0.75),
                 value, 15, bold=True, color=THEME["text_light"])


# ── メイン変換処理 ────────────────────────────────────────────────────────────

def convert(data_path: str, output_path: str,
            title: str = "", subtitle: str = "",
            encoding: str = "utf-8",
            template_path: str | Path | None = None) -> None:
    df = load_data(data_path, encoding)
    total = len(df)

    if not title and "Survey Name" in df.columns:
        val = df["Survey Name"].dropna()
        if len(val):
            title = str(val.iloc[0])
    if not title:
        title = Path(data_path).stem

    col_types = classify_columns(df)

    print(f"ファイル : {data_path}")
    print(f"回答数   : {total} 件 / {len(df.columns)} 列")
    print()
    for col, t in col_types.items():
        icon = {"personal": "⛔", "metadata": "ℹ️",
                "categorical": "📊", "high_cardinality": "📊",
                "text": "📝", "empty": "—", "appendix_company": "📋"}.get(t, "  ")
        print(f"  {icon} [{t:18s}] {col[:55]}")

    path = Path(template_path) if template_path else DEFAULT_TEMPLATE_PATH
    use_template = path.exists()
    if use_template:
        prs = Presentation(str(path))
        # サンプルは「タイトル」「サマリ」の2枚のみ残す
        while len(prs.slides) > 2:
            rId = prs.slides._sldIdLst[-1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[-1]
        print(f"テンプレート: {path}（サンプル2枚のみ残し、タイトル・サマリは新規生成しません）")
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        print("テンプレートなし: デフォルトサイズで新規作成")

    if not use_template:
        add_title_slide(prs, title, subtitle, total)
        add_summary_slide(prs, df, col_types, total)

    q_num = 1
    appendix_company_cols = [c for c, t in col_types.items() if t == "appendix_company"]

    for col, q_type in col_types.items():
        if q_type in ("personal", "metadata", "empty", "appendix_company"):
            continue

        label = str(col)
        col_idx = df.columns.get_loc(col) if col in df.columns else -1

        # 元データのP列（16列目）は円グラフ1枚で表示
        if col_idx == PIE_ONLY_COLUMN_INDEX:
            add_pie_only_slide(prs, df[col], label, q_num, total)
            q_num += 1
            continue

        if q_type in ("categorical", "high_cardinality"):
            s = df[col].dropna().astype(str).str.strip()
            is_4scale = any(v in s.values for v in SCALE_4_ORDER)
            add_categorical_slide(
                prs, df[col], label, q_num, total,
                ordered_values = SCALE_4_ORDER if is_4scale else None,
                colors         = SCALE_4_COLORS if is_4scale else None,
                top_n          = TOP_N_HIGH_CARDINALITY if q_type == "high_cardinality" else None,
            )

        elif q_type == "text":
            add_text_slides(prs, df[col], label, q_num, total)

        q_num += 1

    # Appendix: セクションタイトル + 会社名などを表で最後に追加
    if appendix_company_cols:
        add_section_title_slide(prs, "Appendix")
    for col in appendix_company_cols:
        add_appendix_company_slides(prs, df[col], str(col), total)

    prs.save(output_path)
    print(f"\n✅ 保存完了: {output_path}  ({len(prs.slides)} スライド)")


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description="CSV/XLSXアンケート → PowerPoint変換（個人情報除外・ネイティブチャート）"
    )
    parser.add_argument("file", help="入力ファイル (.csv または .xlsx)")
    parser.add_argument("-o", "--output", help="出力PPTXパス（省略時: 同名.pptx）")
    parser.add_argument("--title",    default="", help="タイトル（省略時: Survey Name列 or ファイル名）")
    parser.add_argument("--subtitle", default="", help="サブタイトル")
    parser.add_argument("--encoding", default="utf-8", help="CSVエンコーディング")
    parser.add_argument("--template", default="", help="テンプレートPPTXパス（省略時: template/design.pptx）")
    args = parser.parse_args()

    input_path = Path(args.file)
    if not input_path.exists():
        print(f"エラー: {input_path} が見つかりません", file=sys.stderr)
        sys.exit(1)

    output_path = args.output or str(input_path.with_suffix(".pptx"))
    template_path = args.template or None
    convert(str(input_path), output_path,
            title=args.title, subtitle=args.subtitle, encoding=args.encoding,
            template_path=template_path)


if __name__ == "__main__":
    main()
